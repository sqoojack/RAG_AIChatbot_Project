
# Upload_file.py
from utils import io, pytesseract, base64, requests, re, json, ffmpeg, tempfile, os, whisper
from langchain.schema import Document
from pptx import Presentation
from docx import Document as DocxDocument
from PyPDF2 import PdfReader
from PIL import Image, ImageFile
from config import ollama_url

# 使用多模態模型去產生圖片描述並回傳 (ex: LLaVA), ollama_url: Ollama的API位置
def get_image_caption(image_bytes, ollama_url, img_model):
    # 將圖片轉成 base64 編碼格式，符合 Ollama 接口格式
    b64_image = base64.b64encode(image_bytes).decode()

    # 建立 POST 請求內容，傳入 prompt 與圖片（以 base64 格式）
    payload = {
        "model": img_model,
        "prompt": "請描述這張圖片的內容，讓視障者可以理解。",
        "images": [b64_image]
    }

    # 發送請求給 Ollama，取得LLM對於prompt以及圖像所生成的回應, json=payload: 自動把 payload 轉成 JSON 字串, 因為這是Ollama API所預期的輸入格式
    response = requests.post(f"{ollama_url}/api/generate", json=payload)

    # 將生成的回應拼接成包含多個json結構的字串
    result = ""
    for line in response.iter_lines():
        if line:
            result += line.decode("utf-8")

    responses = []

    # 使用正則表達式擷取所有 JSON 格式的字串片段
    json_objects = re.findall(r'\{.*?"response":.*?\}', result)

    # 嘗試逐一解析 JSON 片段，並提取 "response" 欄位內容
    for obj in json_objects:
        try:
            data = json.loads(obj)
            if "response" in data:
                responses.append(data["response"])
        except json.JSONDecodeError:
            continue  # 忽略格式錯誤的片段

    # 將所有回傳的片段組合成一段完整圖像描述
    return "".join(responses).strip()

# 載入上傳的檔案，並轉為 Langchain 的 Document 格式, uploaded_file: user上傳的檔案

def load_documents_from_upload(uploaded_file, img_model):

    # 取得檔案名稱與內容
    filename = uploaded_file.name.lower()
    content = uploaded_file.read()

    # === 處理 PDF 檔 ===
    if filename.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(content))  # 用 BytesIO 讀取二進位內容
        docs = []
        for idx, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():  # 若該頁有文字才加入
                docs.append(Document(
                    page_content=text.strip(),
                    metadata={
                        "source": uploaded_file.name,
                        "type": "pdf",
                        "page": idx + 1  # 從 1 開始標記頁碼
                    }
                ))
        return docs

    # === 處理圖檔（jpg, png, jpeg）===
    elif filename.endswith((".jpg", ".jpeg", ".png")):
        ImageFile.LOAD_TRUNCATED_IMAGES = True  # 允許載入不完整圖片
        image = Image.open(io.BytesIO(content)).convert("RGB")  # 開啟圖片並轉為 RGB 模式

        # 使用 Tesseract 執行 OCR（中文＋英文）
        ocr_text = pytesseract.image_to_string(image, lang="chi_tra+eng")
        ocr_lines = [l.strip() for l in ocr_text.splitlines() if l.strip()]
        ocr_paragraph = "\n".join(ocr_lines)

        # 使用圖像模型（如 LLaVA）取得圖片描述
        caption = get_image_caption(content, ollama_url, img_model) or "(無回應)"

        # 組合 OCR + 圖說為完整內容
        full_text = f"""【圖像描述】\n{caption}\n\n【圖片內文字 (OCR)】\n{ocr_paragraph}"""

        return [
            Document(
                page_content=full_text.strip(),
                metadata={
                    "source": uploaded_file.name,
                    "type": "image+ocr+llava"
                }
            )
        ]
    
    elif filename.endswith(".docx"):
        text = ""
        doc = DocxDocument(io.BytesIO(content))
        for para in doc.paragraphs:
            if para.text.strip():
                text += para.text.strip() + "\n"
        return [Document(page_content=text.strip(),
                    metadata={"source": uploaded_file.name, "type": "doc"})]
    
    # 目前不支援doc舊檔案
    elif filename.endswith(".doc"):   
        # 1. 寫入臨時檔案
        with tempfile.TemporaryDirectory() as tmp:
            path_doc = os.path.join(tmp, "input.doc")
            with open(path_doc, "wb") as f:
                f.write(content)

            # 2. 啟動 Word COM 讀取文字
            word = client.Dispatch("Word.Application")
            word.Visible = False
            doc  = word.Documents.Open(path_doc, ReadOnly=True)
            text = doc.Content.Text
            doc.Close(False)
            word.Quit()

        # 3. 回傳 LangChain Document
        return [
            Document(
                page_content=text.strip(),
                metadata={"source": uploaded_file.name, "type": "doc"})]
    
    # 處理ppt檔
    elif filename.endswith((".pptx", ".ppt")):
        prs = Presentation(io.BytesIO(content))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        text += txt + "\n"
        return [Document(page_content=text.strip(),
                        metadata={"source": uploaded_file.name, "type": "ppt"})]
    
    elif filename.endswith(".mp3"):
        with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp:
            tmp.write(content)
            tmp_path = tmp.name
        whisper_model = whisper.load_model("base")
        result = whisper_model.transcribe(tmp_path)    # 用 Whisper 做轉錄
        transcript = result["text"].strip()
        
        return [Document(page_content=transcript, metadata={"source": uploaded_file.name, "type": "audio"})]

    # 處理 mp4 (先只抽音訊, 把他轉成mp3後採用跟mp3一樣的處理方式) 要下載ffmpeg才能轉檔, 目前公司會擋所以暫時無法用mp4
    elif filename.endswith(".mp4"):
        mp3_bytes, _ = (ffmpeg.input("pipe:0").output("pipe:1", format="mp3", acodec="libmp3lame").run(input=content, capture_stdout=True, capture_stderr=True))
        with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as tmp:
            tmp.write(mp3_bytes)
            tmp_path = tmp.name
        whisper_model = whisper.load_model("base")
        result = whisper_model.transcribe(tmp_path)    # 用 Whisper 做轉錄
        transcript = result["text"].strip()
        return [Document(page_content=transcript, metadata={"source": uploaded_file.name, "type": "audio"})]

    # === 處理純文字檔 .txt===
    else:
        try:
            text = content.decode("utf-8")         # 使用 UTF-8 解碼
        except UnicodeDecodeError:
            text = content.decode("big5", errors="ignore")      # 若無法解碼，改用 Big5（台灣常見編碼）

        # 移除空白行並重組內容
        full_text = "\n".join([l.strip() for l in text.splitlines() if l.strip()])

        return [
            Document(
                page_content=full_text,
                metadata={
                    "source": uploaded_file.name,
                    "type": "txt"
                }
            )
        ]


# 處理多個檔案，並呼叫 load_documents_from_upload
def process_uploaded_files(uploaded_files, img_model):
    docs = []
    for f in uploaded_files or []:
        docs.extend(load_documents_from_upload(f, img_model))   # 逐一加入目前的docs中
    return docs